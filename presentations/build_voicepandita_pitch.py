from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = Path(__file__).resolve().parent
LOGO = ROOT / "public" / "icon.jpg"

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(18, 39, 36)
FOREST = RGBColor(5, 104, 78)
MINT = RGBColor(222, 247, 239)
PALE = RGBColor(247, 250, 247)
GOLD = RGBColor(247, 181, 32)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(84, 105, 100)
RED = RGBColor(224, 86, 67)


def rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is not None:
        shape.line.color.rgb = line
    return shape


def text(slide, value, x, y, w, h, size=24, color=INK, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return box


def circle(slide, label, x, y, d, fill, size=20, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill; shape.line.fill.background()
    tf = shape.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return shape


def brand(slide, number):
    slide.shapes.add_picture(str(LOGO), Inches(.48), Inches(.38), Inches(.5), Inches(.5))
    text(slide, "VoicePandita", 1.08, .43, 2.3, .35, 18, INK, True)
    text(slide, f"0{number}", 12.35, .43, .45, .3, 11, MUTED, True, align=PP_ALIGN.RIGHT)


def footer(slide, label):
    rect(slide, .48, 7.12, 12.35, .015, FOREST, False)
    text(slide, label, .48, 7.17, 8, .18, 9, MUTED)
    text(slide, "AI Build Fest • Top 100", 10.4, 7.17, 2.43, .18, 9, FOREST, True, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    # python-pptx 1.0 exposes a writable notes text frame.
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        pass


prs = Presentation()
prs.slide_width, prs.slide_height = W, H
blank = prs.slide_layouts[6]

# Slide 1 — problem / hook
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, PALE, False)
rect(s, 8.95, 0, 4.383, 7.5, FOREST, False)
brand(s, 1)
text(s, "A great tutor should not depend on", .65, 1.28, 7.7, .55, 25, FOREST, True)
text(s, "language, bandwidth,\nor ability.", .65, 1.84, 7.85, 1.75, 47, INK, True)
text(s, "Yet millions of learners are still excluded by English-first tools, unstable internet, and inaccessible formats.", .67, 3.88, 7.55, 1.0, 22, MUTED)
rect(s, .67, 5.23, 7.55, 1.15, MINT)
circle(s, "বাংলা", .91, 5.45, .72, FOREST, 13)
text(s, "VoicePandita brings the tutor to the learner—\nin the way they can learn best.", 1.85, 5.48, 6.0, .72, 20, INK, True)
text(s, "THE GAP", 9.45, 1.05, 2.7, .28, 13, GOLD, True)
for i, (n, title, sub) in enumerate([
    ("1", "Language", "Bangla and indigenous learners"),
    ("2", "Connectivity", "Low-data and offline contexts"),
    ("3", "Accessibility", "Deaf learners and varied needs"),
]):
    yy = 1.62 + i * 1.55
    circle(s, n, 9.43, yy, .55, GOLD, 13, INK)
    text(s, title, 10.2, yy-.02, 2.5, .35, 21, WHITE, True)
    text(s, sub, 10.2, yy+.38, 2.55, .48, 13, RGBColor(213, 239, 230))
footer(s, "1-minute opportunity pitch")
add_notes(s, "0:00–0:17 — In Bangladesh, a student's access to a great tutor still depends on language, bandwidth, and ability. English-first tools, unstable internet, and inaccessible formats leave too many learners behind. VoicePandita was built to close that gap.")

# Slide 2 — product
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, WHITE, False)
brand(s, 2)
text(s, "Ask by voice. Learn with context.", .65, 1.16, 8.8, .7, 35, INK, True)
text(s, "A Bangladesh-first AI learning companion", .67, 1.9, 7.7, .42, 18, FOREST, True)

steps = [
    ("1", "ASK", "Voice • text • photo • PDF"),
    ("2", "GROUND", "Curriculum + retrieval context"),
    ("3", "EXPLAIN", "Bangla • visual • voice • BdSL*"),
    ("4", "CONTINUE", "Offline fallback + learning memory"),
]
for i, (n, title, sub) in enumerate(steps):
    xx = .67 + i * 3.08
    rect(s, xx, 2.75, 2.75, 2.45, MINT if i % 2 == 0 else PALE, line=RGBColor(205, 229, 221))
    circle(s, n, xx+.22, 2.98, .55, FOREST if i < 3 else GOLD, 14, WHITE if i < 3 else INK)
    text(s, title, xx+.22, 3.72, 2.28, .34, 16, FOREST, True)
    text(s, sub, xx+.22, 4.14, 2.28, .72, 16, INK, True)
    if i < 3:
        text(s, "→", xx+2.77, 3.65, .3, .4, 22, FOREST, True, align=PP_ALIGN.CENTER)

rect(s, .67, 5.58, 12.0, .92, FOREST)
text(s, "LIVE PROTOTYPE", .95, 5.86, 1.65, .26, 12, GOLD, True)
text(s, "Next.js PWA  •  Supabase/pgvector  •  AI provider fallbacks  •  privacy-aware design", 2.75, 5.79, 9.45, .4, 17, WHITE, True)
text(s, "* BdSL text-to-sign playback is a prototype; low-resource language output uses verified-first safeguards.", .7, 6.67, 10.9, .25, 10, MUTED)
footer(s, "Implemented in current branch; claims scoped to prototype status")
add_notes(s, "0:17–0:42 — Our live prototype lets a student ask by voice, text, photo, or PDF. It grounds answers in curriculum context, then explains in Bangla through text, voice, diagrams, and a prototype BdSL avatar. Offline packs and deterministic fallbacks keep learning available when connectivity or AI providers fail.")

# Slide 3 — impact and opportunity
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.333, 7.5, PALE, False)
brand(s, 3)
text(s, "From prototype to", .65, 1.15, 7.5, .55, 31, INK, True)
text(s, "inclusive learning infrastructure.", .65, 1.7, 7.85, .72, 34, FOREST, True)

for i, (title, desc, color) in enumerate([
    ("RESEARCH", "Validate learning gains, low-resource language safety, and accessible AI tutoring.", FOREST),
    ("PILOTS", "Co-design with schools and communities in low-connectivity settings.", RGBColor(16, 132, 130)),
    ("SCALE", "Build a sustainable Bangladesh launch and adapt the platform for global underserved markets.", RGBColor(161, 113, 9)),
]):
    yy = 2.77 + i * 1.05
    rect(s, .67, yy, 7.35, .84, WHITE, line=RGBColor(211, 226, 220))
    rect(s, .67, yy, .1, .84, color, False)
    text(s, title, .95, yy+.17, 1.45, .3, 14, color, True)
    text(s, desc, 2.35, yy+.13, 5.3, .5, 15, INK, True)

rect(s, 8.62, 1.22, 4.05, 5.28, FOREST)
text(s, "OUR ASK", 9.05, 1.72, 2.0, .35, 14, GOLD, True)
text(s, "Help us prove,\npilot, and scale.", 9.05, 2.2, 3.1, 1.25, 31, WHITE, True)
text(s, "Research partners\nSchool pilot access\nIncubation & mentorship\nMarket-entry guidance", 9.05, 3.74, 3.05, 1.55, 18, RGBColor(222, 247, 239), True)
rect(s, 9.05, 5.68, 2.95, .5, GOLD)
text(s, "voice-pandita.vercel.app", 9.13, 5.82, 2.8, .2, 11, INK, True, align=PP_ALIGN.CENTER)
text(s, "Every learner deserves a tutor that speaks their language—and understands their reality.", .67, 6.36, 7.55, .55, 17, INK, True)
footer(s, "Opportunity: research • incubation • commercialization • global expansion")
add_notes(s, "0:42–1:00 — We now want to move from a Top 100 hackathon prototype to validated inclusive learning infrastructure. We are seeking research partners, school pilots, incubation, and market-entry guidance. Our vision is simple: every learner deserves a tutor that speaks their language—and understands their reality. We would love to build that next stage with you.")

prs.core_properties.title = "VoicePandita — 1 Minute Opportunity Pitch"
prs.core_properties.subject = "Beyond the Build webinar pitch"
prs.core_properties.author = "VoicePandita Team"
prs.save(OUT / "VoicePandita_1_Minute_Pitch.pptx")
print(OUT / "VoicePandita_1_Minute_Pitch.pptx")
